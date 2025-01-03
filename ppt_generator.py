from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import re
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import numpy as np

def extract_metrics(text):
    """Extract numerical data and metrics from analysis text."""
    metrics = {
        'values': [],
        'labels': [],
        'percentages': []
    }
    
    try:
        # Extract metrics with format "MetricName: Value"
        metric_matches = re.findall(r'([^:\n]+):\s*(\d+(?:\.\d+)?)', text)
        for label, value in metric_matches:
            if '%' in label or '%' in value:
                metrics['percentages'].append(True)
            else:
                metrics['percentages'].append(False)
            metrics['values'].append(float(value.strip('%')))
            metrics['labels'].append(label.strip())
        
        return metrics
    except Exception as e:
        print(f"Error extracting metrics: {str(e)}")
        return {'values': [], 'labels': [], 'percentages': []}

def add_chart_slide(prs, metrics, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    """Add a slide with the specified chart type."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add chart
    chart_data = CategoryChartData()
    chart_data.categories = metrics['labels']
    chart_data.add_series('Values', metrics['values'])
    
    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5.5)
    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
    
    # Customize chart
    chart.has_legend = True
    chart.has_title = True
    chart.chart_title.text_frame.text = "Performance Metrics"

def create_presentation(analysis):
    """Create presentation with standard formatting"""
    try:
        prs = Presentation()
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        title.text = "Player Performance Analysis"
        if subtitle:
            subtitle.text = "Performance Analysis Report"
        
        # Analysis sections
        sections = analysis.split('\n\n')
        for section in sections:
            if section.strip():
                # Create content slide for each major section
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # Extract section title and content
                lines = section.strip().split('\n')
                if lines:
                    title = slide.shapes.title
                    title.text = lines[0].strip()
                    
                    if len(lines) > 1:
                        body_shape = slide.shapes.placeholders[1]
                        tf = body_shape.text_frame
                        
                        for line in lines[1:]:
                            p = tf.add_paragraph()
                            p.text = line.strip()
                            p.level = 1 if line.startswith('-') else 0
        
        # Extract and add metrics charts
        metrics = extract_metrics(analysis)
        if metrics['values']:
            _add_metrics_dashboard(prs, metrics)
            
            # Add hole analysis charts
            hole_pattern = r'Hole \d+:'
            if re.search(hole_pattern, analysis):
                hole_stats = {}
                for match in re.finditer(r'Hole (\d+):\s*(\d+(?:\.\d+)?)\s*\(Handicap Group Avg:\s*(\d+(?:\.\d+)?),\s*Field Avg:\s*(\d+(?:\.\d+)?)\)', analysis):
                    hole_num = int(match.group(1))
                    hole_stats[f'Hole {hole_num}'] = {
                        'player_score': float(match.group(2)),
                        'handicap_group_avg': float(match.group(3)),
                        'field_average': float(match.group(4))
                    }
                if hole_stats:
                    _add_hole_analysis_slide(prs, hole_stats)
        
        # Add hole analysis if available
        hole_pattern = r'Hole \d+.*?player_score: (\d+\.?\d*)'
        if re.search(hole_pattern, analysis):
            _add_hole_analysis_slide(prs, metrics)
            
        # Add time analysis if available
        if 'morning_avg' in metrics and 'afternoon_avg' in metrics:
            _add_time_analysis_slide(prs, metrics)
        
        # Save presentation
        temp_ppt = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_ppt.name)
        return temp_ppt.name
        
    except Exception as e:
        raise Exception(f"Error creating presentation: {str(e)}")

def _structure_content(analysis):
    """Better structure the content for presentation"""
    sections = []
    current_section = []
    
    for line in analysis.split('\n'):
        if line.strip():
            if any(key in line.lower() for key in ['summary:', 'overview:', 'metrics:', 'conclusion:']):
                if current_section:
                    sections.append('\n'.join(current_section))
                current_section = [line]
            else:
                current_section.append(line)
    
    if current_section:
        sections.append('\n'.join(current_section))
    
    return sections

def _add_content_slide(prs, content, styles, layouts):
    """Add content slide with template styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Add title
    if hasattr(slide.shapes, "title"):
        title = slide.shapes.title
        first_line = content.split('\n')[0]
        title.text = first_line
        
        # Apply title styling
        if 'content_title' in styles:
            title_style = styles['content_title']
            if hasattr(title.text_frame.paragraphs[0], 'font'):
                font = title.text_frame.paragraphs[0].font
                font.name = title_style.get('font_name', 'Calibri')
                font.size = title_style.get('size', Pt(32))
                font.bold = title_style.get('bold', True)

    # Add content
    if hasattr(slide.shapes, "placeholders"):
        body_shape = slide.shapes.placeholders[1]
        if body_shape.has_text_frame:
            tf = body_shape.text_frame
            tf.text = '\n'.join(content.split('\n')[1:])

def _apply_template_style(shape, style):
    """Apply template styling to a shape"""
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            if hasattr(paragraph, "font"):
                font = paragraph.font
                font.name = style.get('font_name', 'Calibri')
                font.size = style.get('size', Pt(32))
                font.bold = style.get('bold', True)

def _add_summary_slide(prs, analysis):
    """Add summary slide with basic styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    # Extract summary section
    summary = re.search(r'summary:?(.+?)(?=\n\n|\Z)', analysis, re.I | re.S)
    if summary:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = summary.group(1).strip()
        
        # Apply basic bullet points
        for paragraph in tf.paragraphs[1:]:
            paragraph.level = 1

def _add_chart_slides(prs, metrics):
    """Add chart slides based on metrics and template preferences."""
    if not metrics['values']:
        return

    # Add column chart for general metrics
    add_chart_slide(prs, metrics)
    
    # If we have percentage metrics, create a pie chart
    percentage_metrics = {
        'values': [v for v, is_pct in zip(metrics['values'], metrics['percentages']) if is_pct],
        'labels': [l for l, is_pct in zip(metrics['labels'], metrics['percentages']) if is_pct],
        'percentages': [True] * len([p for p in metrics['percentages'] if p])
    }
    
    if percentage_metrics['values']:
        add_chart_slide(prs, percentage_metrics, XL_CHART_TYPE.PIE)

def _add_metrics_dashboard(prs, metrics):
    """Add a dashboard-style slide with charts"""
    if not metrics['values']:
        return
        
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using blank layout
        title = slide.shapes.title
        title.text = "Performance Metrics"
        
        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = metrics['labels']
        chart_data.add_series('Values', metrics['values'])
        
        # Add chart to slide
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # Customize chart
        chart.has_legend = True
        chart.has_title = True
        chart.chart_title.text_frame.text = "Performance Metrics Overview"
        
    except Exception as e:
        print(f"Error adding metrics dashboard: {str(e)}")

def _add_trend_analysis(prs, analysis, metrics):
    """Add trend analysis with line charts"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "Trend Analysis"
    
    # Extract trend data
    trend_matches = re.findall(r'trend:?\s*([^:]+):\s*(\d+(?:\.\d+)?)', analysis, re.I)
    if trend_matches:
        trend_data = {label.strip(): float(value) for label, value in trend_matches}
        _add_small_chart(slide, trend_data, XL_CHART_TYPE.LINE,
                        Inches(1), Inches(1.5), Inches(8), Inches(5))

def _add_small_chart(slide, data, chart_type, x, y, cx, cy):
    """Add a chart to the slide at specified position"""
    if not data:
        return  # Ensure data is not empty
    
    chart_data = CategoryChartData()
    chart_data.categories = list(data.keys())
    chart_data.add_series('Values', list(data.values()))
    
    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
    chart.has_legend = True

def _add_recommendations_slide(prs, analysis):
    """Add recommendations with impact analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Recommendations & Impact Analysis"
    
    # Extract recommendations and their impacts
    recommendations = re.findall(r'recommend.*?:\s*([^(]+)\((\d+(?:\.\d+)?%?)\)', analysis, re.I)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for rec, impact in recommendations:
        p = tf.add_paragraph()
        p.text = f"• {rec.strip()} (Impact: {impact})"
        p.level = 0

def _add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    
    title.text = "Document Analysis Report"
    if subtitle:
        subtitle.text = "AI-Powered Analysis"

def _add_analysis_slide(prs, section):
    """Add analysis slide with basic formatting"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Using two-content layout
    
    # Extract section title and content
    title = slide.shapes.title
    section_lines = section.split('\n')
    title.text = section_lines[0]
    
    # Extract metrics and key points
    metrics = []
    key_points = []
    
    for line in section_lines[1:]:
        if ':' in line and any(char.isdigit() for char in line):
            metrics.append(line.strip())
        elif line.strip():
            key_points.append(line.strip())
    
    # Add content to left and right placeholders
    if len(slide.placeholders) > 2:
        left_content = slide.placeholders[1]
        right_content = slide.placeholders[2]
        
        # Add metrics to left side
        if metrics:
            tf = left_content.text_frame
            tf.text = "Key Metrics"
            for metric in metrics:
                p = tf.add_paragraph()
                p.text = f"• {metric}"
                p.level = 1
        
        # Add key points to right side
        if key_points:
            tf = right_content.text_frame
            tf.text = "Analysis"
            for point in key_points:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.level = 1

def _add_segment_analysis(prs, analysis):
    """Add segment analysis slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Segment Analysis"
    
    # Extract segment data
    segments = re.findall(r'segment:?\s*([^:]+):\s*(\d+(?:\.\d+)?)', analysis, re.I)
    if segments:
        segment_data = {label.strip(): float(value) for label, value in segments}
        _add_small_chart(slide, segment_data, XL_CHART_TYPE.BAR_CLUSTERED,
                        Inches(1), Inches(1.5), Inches(8), Inches(5))

def _add_hole_analysis_slide(prs, hole_stats):
    """Add hole-by-hole comparison charts"""
    # Create first slide for line chart comparison
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide1.shapes.title
    title.text = "Hole-by-Hole Performance"
    
    # Prepare data
    holes = sorted([h for h in hole_stats.keys() if h.startswith('Hole')])
    
    # Line chart for score comparison
    chart_data = CategoryChartData()
    chart_data.categories = holes
    
    player_scores = [hole_stats[hole]['player_score'] for hole in holes]
    field_averages = [hole_stats[hole]['field_average'] for hole in holes]
    handicap_averages = [hole_stats[hole].get('handicap_group_avg', 0) for hole in holes]
    
    chart_data.add_series('Player Score', player_scores)
    chart_data.add_series('Field Average', field_averages)
    chart_data.add_series('Handicap Group', handicap_averages)
    
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
    chart = slide1.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.has_title = True
    chart.chart_title.text_frame.text = "Score Comparison by Hole"

    # Create second slide for performance difference
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide2.shapes.title
    title.text = "Hole Performance Analysis"
    
    # Calculate differences from average
    diff_data = CategoryChartData()
    diff_data.categories = holes
    
    # Calculate differences (negative is better in golf)
    field_differences = [-1 * (ps - fa) for ps, fa in zip(player_scores, field_averages)]
    handicap_differences = [-1 * (ps - ha) for ps, ha in zip(player_scores, handicap_averages)]
    
    diff_data.add_series('vs Field Average', field_differences)
    diff_data.add_series('vs Handicap Group', handicap_differences)
    
    # Add bar chart showing differences
    chart = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, diff_data
    ).chart
    
    chart.has_legend = True
    chart.has_title = True
    chart.chart_title.text_frame.text = "Performance vs Averages (Negative is Better)"

def _add_time_analysis_slide(prs, time_stats):
    """Add time of day analysis chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Morning vs Afternoon Performance"
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Morning', 'Afternoon']
    chart_data.add_series('Average Score', [
        time_stats['morning_avg'],
        time_stats['afternoon_avg']
    ])
    
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.has_title = True
