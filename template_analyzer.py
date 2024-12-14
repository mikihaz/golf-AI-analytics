from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt
import json
import re

class TemplateAnalyzer:
    def __init__(self):
        self.style_patterns = {
            'fonts': {},
            'colors': {},
            'layouts': [],
            'chart_types': [],
            'spacing': {},
            'section_patterns': [],
            'bullet_patterns': [],
            'slide_count': 0
        }

    def learn_from_template(self, template_path):
        """Extract detailed styling and structure from template"""
        prs = Presentation(template_path)
        template_data = {
            'styles': self._extract_styles(prs),
            'layouts': self._analyze_layouts(prs),
            'charts': self._analyze_charts(prs),
            'color_scheme': self._extract_color_scheme(prs),
            'structure': self._analyze_structure(prs),
            'formatting': self._analyze_formatting(prs)
        }
        return template_data

    def _extract_styles(self, prs):
        """Extract detailed styling information"""
        styles = {
            'title': {'font_name': 'Calibri', 'size': Pt(44), 'bold': True},
            'content_title': {'font_name': 'Calibri', 'size': Pt(32), 'bold': True},
            'body': {'font_name': 'Calibri', 'size': Pt(18), 'bold': False}
        }

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if hasattr(paragraph, "font"):
                            font = paragraph.font
                            if shape.name == "Title":
                                styles['title'].update({
                                    'font_name': font.name or 'Calibri',
                                    'size': font.size or Pt(44),
                                    'bold': font.bold
                                })
        return styles

    def _analyze_layouts(self, prs):
        """Enhanced layout analysis"""
        layouts = []
        for slide in prs.slides:
            layout_info = {
                'type': self._determine_detailed_type(slide),
                'placeholders': self._analyze_placeholders(slide),
                'background': self._capture_background(slide),
                'shapes': self._analyze_shapes(slide)
            }
            layouts.append(layout_info)
        return layouts

    def _determine_detailed_type(self, slide):
        """Determine detailed slide type based on content and structure"""
        if not slide.shapes:
            return 'blank'
        
        # Check for title slide
        if len(slide.shapes) <= 2 and hasattr(slide.shapes.title, "text"):
            return 'title'
            
        # Check for charts
        charts = [shape for shape in slide.shapes if hasattr(shape, "chart")]
        if charts:
            if len(charts) > 1:
                return 'dashboard'
            return 'chart'
            
        # Check content structure
        placeholders = [shape for shape in slide.shapes if shape.is_placeholder]
        if len(placeholders) > 1:
            return 'two_content'
            
        return 'content'

    def _analyze_placeholders(self, slide):
        """Analyze placeholder structure and formatting"""
        placeholders = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholders.append({
                    'idx': shape.placeholder_format.idx,
                    'type': shape.placeholder_format.type if hasattr(shape.placeholder_format, 'type') else None,
                    'position': (shape.left, shape.top, shape.width, shape.height),
                    'has_text': hasattr(shape, "text_frame")
                })
        return placeholders

    def _capture_background(self, slide):
        """Capture background styling"""
        try:
            if slide.background.fill.type:
                return {
                    'type': str(slide.background.fill.type),
                    'color': slide.background.fill.fore_color.rgb if hasattr(slide.background.fill, 'fore_color') else None
                }
        except:
            pass
        return None

    def _analyze_shapes(self, slide):
        """Analyze shape properties and styling"""
        shapes = []
        for shape in slide.shapes:
            if not shape.is_placeholder:
                shape_info = {
                    'type': shape.shape_type,
                    'position': (shape.left, shape.top, shape.width, shape.height),
                    'has_text': hasattr(shape, "text_frame"),
                    'has_chart': hasattr(shape, "chart")
                }
                shapes.append(shape_info)
        return shapes

    def _analyze_charts(self, prs):
        """Analyze chart types and their usage patterns"""
        chart_patterns = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "chart"):
                    chart_info = {
                        'type': str(shape.chart.chart_type),
                        'has_legend': shape.chart.has_legend,
                        'style': shape.chart.style,
                    }
                    chart_patterns.append(chart_info)
        return chart_patterns

    def _extract_color_scheme(self, prs):
        """Extract color scheme from template"""
        colors = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "fill"):
                    if shape.fill.type:
                        try:
                            colors.add(shape.fill.fore_color.rgb)
                        except:
                            pass
        return list(colors)

    def _analyze_structure(self, prs):
        """Analyze the presentation structure and patterns"""
        structure = {
            'sections': [],
            'headings': [],
            'bullet_patterns': [],
            'slide_sequence': []
        }
        
        for slide in prs.slides:
            slide_type = self._determine_slide_type(slide)
            structure['slide_sequence'].append(slide_type)
            
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    if text:
                        # Capture heading patterns
                        if shape.name == "Title":
                            structure['headings'].append(text)
                        
                        # Capture bullet point patterns
                        if shape.text_frame.paragraphs[0].level:
                            structure['bullet_patterns'].append(
                                self._extract_bullet_pattern(shape.text_frame)
                            )
        
        return structure

    def _determine_slide_type(self, slide):
        """Determine the type and purpose of a slide"""
        if not slide.shapes:
            return 'blank'
        
        has_chart = any(hasattr(shape, "chart") for shape in slide.shapes)
        if has_chart:
            return 'chart'
            
        text_content = ' '.join(
            shape.text_frame.text for shape in slide.shapes 
            if hasattr(shape, "text_frame")
        ).lower()
        
        if any(key in text_content for key in ['summary', 'overview']):
            return 'summary'
        elif any(key in text_content for key in ['conclusion', 'next steps']):
            return 'conclusion'
        else:
            return 'content'

    def generate_prompt(self, template_data):
        """Generate analysis prompt matching template structure"""
        if not template_data:
            return self._default_prompt()

        structure = template_data.get('structure', {})
        sections = structure.get('sections', [])
        
        prompt = "Analyze the content and structure it exactly as follows:\n\n"
        
        # Match section headings from template
        for heading in structure.get('headings', []):
            if heading.strip():
                prompt += f"# {heading}\n"
        
        # Match bullet patterns
        bullet_patterns = structure.get('bullet_patterns', [])
        if bullet_patterns:
            prompt += "\nUse similar bullet point structure as:\n"
            for pattern in bullet_patterns[:2]:  # Show a couple examples
                prompt += f"- {pattern}\n"
        
        prompt += "\nEnsure numerical data is presented in a format matching the template charts.\n"
        return prompt

    def _extract_bullet_pattern(self, text_frame):
        """Extract bullet point pattern and structure"""
        pattern = []
        for paragraph in text_frame.paragraphs:
            if paragraph.level == 0:
                pattern.append('•')
            elif paragraph.level == 1:
                pattern.append('   -')
        return ' '.join(pattern)

    def save_patterns(self, filepath):
        """Save learned patterns to file"""
        with open(filepath, 'w') as f:
            json.dump(self.style_patterns, f)

    def load_patterns(self, filepath):
        """Load previously learned patterns"""
        with open(filepath, 'r') as f:
            self.style_patterns = json.load(f)

    def _analyze_formatting(self, prs):
        """Analyze text formatting patterns in the template"""
        formatting = {
            'paragraph_spacing': {},
            'text_alignment': {},
            'indent_levels': set(),
            'list_styles': [],
            'text_boxes': []
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Analyze paragraph formatting
                    for paragraph in shape.text_frame.paragraphs:
                        if hasattr(paragraph, "alignment"):
                            formatting['text_alignment'][paragraph.alignment] = \
                                formatting['text_alignment'].get(paragraph.alignment, 0) + 1
                        
                        if hasattr(paragraph, "level"):
                            formatting['indent_levels'].add(paragraph.level)
                        
                        # Check for bullet points
                        if hasattr(paragraph, "bullet"):
                            if paragraph.bullet:
                                formatting['list_styles'].append({
                                    'level': paragraph.level,
                                    'character': paragraph.bullet.char if hasattr(paragraph.bullet, 'char') else '•'
                                })
        
        # Convert sets to lists for JSON serialization
        formatting['indent_levels'] = list(formatting['indent_levels'])
        return formatting
